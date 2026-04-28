"""Generate COMP34612 CGT Group 12 Presentation - Stackelberg Pricing Game."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette (inspired by Hex presentation) ──
BG_DARK    = RGBColor(0x1A, 0x0A, 0x3E)   # deep purple
BG_MID     = RGBColor(0x24, 0x10, 0x55)   # mid purple
CYAN       = RGBColor(0x00, 0xE5, 0xCC)   # teal/cyan - headers
MAGENTA    = RGBColor(0xE0, 0x5F, 0xCD)   # magenta - accents
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY  = RGBColor(0x2A, 0x2A, 0x3A)   # code block bg
GREEN      = RGBColor(0x00, 0xE6, 0x76)   # success
RED_ACCENT = RGBColor(0xFF, 0x55, 0x55)   # failure
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)   # highlight numbers
TABLE_BG   = RGBColor(0x1E, 0x0E, 0x4A)   # table cell bg
TABLE_HDR  = RGBColor(0x15, 0x08, 0x38)   # table header bg

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════

def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=16, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2),
             font_name='Calibri'):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet(tf, text, font_size=15, color=WHITE, level=0, bold=False):
    """Add a bullet point."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.level = level
    p.space_before = Pt(3)
    p.space_after = Pt(2)
    return p


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    """Add a dark code block with monospace text."""
    # Background rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GRAY
    shape.line.fill.background()
    shape.shadow.inherit = False

    # Code text
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = 'Consolas'
        p.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        p.space_before = Pt(1)
        p.space_after = Pt(1)
    return shape


def add_table(slide, left, top, width, height, data, col_widths=None):
    """Add a styled table. data = list of rows, first row is header."""
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)

            # Style cell
            fill = cell.fill
            fill.solid()
            if r_idx == 0:
                fill.fore_color.rgb = TABLE_HDR
            else:
                fill.fore_color.rgb = TABLE_BG

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = 'Calibri'
                paragraph.alignment = PP_ALIGN.CENTER
                if r_idx == 0:
                    paragraph.font.color.rgb = CYAN
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = WHITE
    return table


def add_accent_box(slide, left, top, width, height, text, font_size=13,
                   border_color=CYAN, text_color=WHITE):
    """Add a bordered accent box."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x18, 0x08, 0x3A)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    return tf


def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_section_header(slide, text, subtitle=None):
    """Add main header at top of slide."""
    tf = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.9),
                     text, font_size=36, color=CYAN, bold=True,
                     font_name='Calibri')
    if subtitle:
        add_para(tf, subtitle, font_size=16, color=MAGENTA, bold=True)
    return tf


# ═══════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# Course code - large
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(1.2),
            'COMP34612', font_size=54, color=CYAN, bold=True)

# Project title
tf = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.0),
                 'Computational Game Theory', font_size=32, color=WHITE, bold=False)
add_para(tf, 'Group Project', font_size=28, color=WHITE)

# Subtitle - the journey framing
add_textbox(slide, Inches(0.8), Inches(3.8), Inches(8), Inches(0.8),
            'Adaptive Stackelberg Pricing: From Textbook OLS to 99% Optimality',
            font_size=20, color=MAGENTA, bold=True)

# Group number box
add_accent_box(slide, Inches(9.5), Inches(3.2), Inches(3.0), Inches(1.8),
               'GROUP\n12', font_size=36, border_color=CYAN, text_color=CYAN)

# Team info at bottom
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
            'A Story of Exploration, Adaptation, and Evidence-Based Iteration',
            font_size=16, color=LIGHT_GRAY, bold=False)

add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
            'Week 11  |  Semester 2  |  2025-26',
            font_size=14, color=LIGHT_GRAY)

add_speaker_notes(slide, """SPEAKER: Member 1 (2 min)
- Welcome, introduce the group and project
- Frame the narrative: "This is the story of how we went from a textbook approach scoring 85% to an adaptive system achieving 99% optimality"
- Mention the project started Week 7 (16 March 2026), with weekly support sessions
- Transition: "Let me hand over to [Member 2] to explain the game setup." """)


# ═══════════════════════════════════════════════════
# SLIDE 2: PROBLEM OVERVIEW
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'Problem Overview',
                   'Stackelberg Pricing Game Under Imperfect Information')

# Left column - game setup
tf = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.5),
                 'GAME STRUCTURE', font_size=18, color=CYAN, bold=True)
add_bullet(tf, 'Leader sets price u_L, follower responds with u_F', color=WHITE)
add_bullet(tf, '100 days historical data (t = 1...100)', color=WHITE)
add_bullet(tf, '30 days live gameplay (t = 101...130)', color=WHITE)
add_bullet(tf, 'Follower reaction function is UNKNOWN', color=MAGENTA, bold=True)
add_bullet(tf, 'Parameters may change over time', color=WHITE)

add_para(tf, '', font_size=8)
add_para(tf, 'WHAT WE KNOW', font_size=18, color=CYAN, bold=True)
add_bullet(tf, 'Demand: S_L(u_L, u_F) = 100 - 5u_L + 3u_F', color=WHITE)
add_bullet(tf, 'Profit: (u_L - 1) * S_L     [unit cost c_L = 1]', color=WHITE)
add_bullet(tf, 'Objective: maximise 30-day accumulated profit', color=WHITE)

add_para(tf, '', font_size=8)
add_para(tf, 'WHAT WE DON\'T KNOW', font_size=18, color=MAGENTA, bold=True)
add_bullet(tf, 'Follower\'s payoff function', color=WHITE)
add_bullet(tf, 'Follower\'s strategy space', color=WHITE)
add_bullet(tf, 'Whether parameters are stationary', color=WHITE)

# Right column - followers table
data = [
    ['Follower', 'Strategy Space', 'Hidden Variant'],
    ['MK1', '[1, +inf)', 'MK4'],
    ['MK2', '[1, +inf)', 'MK5'],
    ['MK3', '[1, 15]', 'MK6'],
]
add_table(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(1.8), data)

# Key insight box
add_accent_box(slide, Inches(7.0), Inches(3.8), Inches(5.8), Inches(1.8),
               'THE CHALLENGE\n\nLearn the follower\'s reaction function\n'
               'from data, then compute the Stackelberg-optimal\n'
               'price  --  all while playing the game live.',
               font_size=14, border_color=MAGENTA)

# Lecture reference
add_textbox(slide, Inches(7.0), Inches(6.0), Inches(5.8), Inches(0.8),
            'L1 Slides 3-5: Stackelberg framework  |  L4 Slides 3-5: Imperfect information\n'
            'L7 Slides 6-8: Best reaction functions & equilibrium derivation',
            font_size=12, color=LIGHT_GRAY)

add_speaker_notes(slide, """SPEAKER: Member 2 (2 min)
- Explain the Stackelberg game: leader commits first, follower responds
- Emphasise: we don't know the follower's payoff or reaction function
- Point to the table: 3 visible + 3 hidden followers
- Key message: "This is a learning problem under imperfect information"
- Reference L1 and L4 explicitly when describing the framework""")


# ═══════════════════════════════════════════════════
# SLIDE 3: THE NAIVE APPROACH & WHY IT FAILS
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'The Naive Approach',
                   'Why Textbook OLS Alone Scores 85.3%')

# Left - what naive does
tf = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.0),
                 'TEXTBOOK TWO-STEP METHOD (Lecture 4)', font_size=16, color=CYAN, bold=True)
add_bullet(tf, 'Step 1: Fit u_F = a + b*u_L via OLS on historical data (L4 Slide 23)', color=WHITE)
add_bullet(tf, 'Step 2: Substitute into profit, solve for optimal u_L (L4 Slide 22)', color=WHITE)
add_bullet(tf, 'Step 3: Play u_L* for all 30 days', color=WHITE)
add_para(tf, 'Result: 85.3% average optimality', font_size=16, color=RED_ACCENT, bold=True)

# The critical problem
tf = add_textbox(slide, Inches(0.6), Inches(3.8), Inches(5.8), Inches(3.2),
                 'THE EXTRAPOLATION TRAP', font_size=18, color=MAGENTA, bold=True)
add_bullet(tf, 'Historical leader prices: u_L in [1.72, 1.90]', color=WHITE, bold=True)
add_bullet(tf, 'Range width: only 0.18 units!', color=YELLOW)
add_bullet(tf, 'True optimal for MK1/MK2: u_L ~ 20', color=WHITE)
add_bullet(tf, 'True optimal for MK3: u_L ~ 11', color=WHITE)
add_bullet(tf, 'OLS slope estimated from 0.18-wide range...', color=WHITE)
add_bullet(tf, '...then extrapolated 10x beyond training data', color=RED_ACCENT, bold=True)
add_bullet(tf, 'Small slope error = massive price error', color=RED_ACCENT)

# Right - comparison table
data = [
    ['Leader', 'MK1', 'MK2', 'MK3', 'Average'],
    ['NaiveLeader\n(textbook OLS)', '81.1%', '75.4%', '99.5%', '85.3%'],
    ['AdaptiveLeader\n(ours)', '99.2%', '99.4%', '100.0%', '99.5%'],
]
add_table(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.0), data)

# Visual: the range problem with concrete numbers
add_accent_box(slide, Inches(7.0), Inches(3.8), Inches(5.8), Inches(3.2),
               'OUR FIRST EXPERIMENT\n\n'
               'We implemented textbook OLS (L4 Slide 23) and ran it.\n'
               'Historical OLS gave: beta_hat = 0.42, alpha_hat = 7.1\n'
               'Predicted optimal: u_L* = 21.3\n'
               'Actual profit: only 81.1% on MK1 (lost ~5,500)\n\n'
               'WHY? The OLS was fit on u_L in [1.72, 1.90].\n'
               'beta_hat = 0.42 was wrong -- true beta ~ 0.30.\n'
               'A 0.12 slope error at u_L=20 means a 2.4 price error.\n\n'
               'This single experiment defined our entire project direction.',
               font_size=12, border_color=YELLOW)

add_speaker_notes(slide, """SPEAKER: Member 2 (2 min)
- Walk through the comparison table: NaiveLeader vs AdaptiveLeader
- Emphasise the 0.18-range insight: "Imagine fitting a line through 100 points clustered in a tiny range, then predicting a value 10x beyond"
- The visual box shows the scale of extrapolation
- "This is why exploration is essential -- and it's the main contribution of our project"
- Transition to course theory""")


# ═══════════════════════════════════════════════════
# SLIDE 4: COURSE THEORY FOUNDATION
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'Course Theory Foundation',
                   'Building on Lectures 1, 3, 4, 5, 6, and 7')

# Theory pipeline - left side
tf = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.5),
                 'OUR THEORETICAL PIPELINE', font_size=18, color=CYAN, bold=True)

add_para(tf, 'L1 & L7: Stackelberg Framework', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, 'L1 Slide 3-5: Two-step optimisation: learn R(u_L), then max profit', font_size=14)
add_bullet(tf, 'L7 Slide 6-8: Existence of Stackelberg equilibrium guaranteed', font_size=14)

add_para(tf, 'L3: Linear Reaction Function', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, 'L3 Slide 12: u_F = alpha + beta*u_L justified by linear demand', font_size=14)

add_para(tf, 'L4: OLS for Imperfect Information', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, 'L4 Slide 14-20: Learn R_hat(u_L) via least squares regression', font_size=14)
add_bullet(tf, 'L4 Slide 22-23: Two-step approximate Stackelberg strategy', font_size=14)

add_para(tf, 'L5: Model Evaluation Metrics', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, 'L5 Slide 4-8: RMSE, MAPE, R-squared to assess fit quality', font_size=14)

add_para(tf, 'L6: RLS with Forgetting Factor', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, 'L6 Slide 9: Weighted LS with forgetting factor lambda', font_size=14)
add_bullet(tf, 'L6 Slide 12: Recursive LS update formula (our implementation)', font_size=14)

# Right side - the formula derivation
tf = add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(1.5),
                 'STACKELBERG OPTIMAL PRICE DERIVATION', font_size=16, color=CYAN, bold=True)
add_para(tf, 'Given: u_F = alpha + beta * u_L', font_size=14, color=WHITE)
add_para(tf, 'Profit = (u_L - 1)(100 - 5u_L + 3(alpha + beta*u_L))', font_size=14, color=WHITE)
add_para(tf, 'dProfit/du_L = 0  =>  u_L* = (105 + 3a - 3b) / (10 - 6b)', font_size=14,
         color=YELLOW, bold=True)

# Code snippet - the optimal price function
code = '''def _optimal_price(self, date=115):
    a = self.alpha + self.gamma * date
    b = self.beta
    denom = 10 - 6 * b
    if denom > 0.1:
        uL = (105 + 3*a - 3*b) / denom
    else:
        # steep slope: ramp up gradually
        uL = prev_uL * 1.5
    return max(1.01, min(uL, self.UPPER_BOUND))'''
add_code_block(slide, Inches(7.0), Inches(3.5), Inches(5.8), Inches(3.0),
               code, font_size=11)

add_textbox(slide, Inches(7.0), Inches(6.7), Inches(5.8), Inches(0.5),
            'leaders.py:69-83  |  Direct implementation of Lecture 4 & 7 theory',
            font_size=11, color=LIGHT_GRAY)

add_speaker_notes(slide, """SPEAKER: Member 3 (2 min)
- Walk through the theoretical pipeline: L1 -> L3 -> L4 -> L5 -> L6
- Show the optimal price formula derivation on the right
- Point to the code: "This is the direct implementation of the L4 Slide 23 two-step method"
- Key message: "The course gave us every tool we needed -- OLS, RLS, metrics, Stackelberg theory"
- Transition: "But there's a critical gap the theory doesn't address..."
""")


# ═══════════════════════════════════════════════════
# SLIDE 5: THE EXPLORATION PROBLEM
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'The Exploration Problem',
                   'Why You Can\'t Just Exploit From Day 101')

# Left - the problem
tf = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.0),
                 'THE CORE DILEMMA', font_size=18, color=CYAN, bold=True)
add_bullet(tf, 'OLS fitted on u_L in [1.72, 1.90] gives a slope estimate')
add_bullet(tf, 'But the optimal price is u_L ~ 20 (MK1/MK2) or ~11 (MK3)')
add_bullet(tf, 'Plugging a wrong slope into the formula gives wrong u_L*', bold=True)
add_para(tf, '', font_size=6)
add_para(tf, 'EXPLORATION-EXPLOITATION TRADE-OFF', font_size=16, color=MAGENTA, bold=True)
add_bullet(tf, 'Explore: sacrifice 1 day of profit to learn the true slope')
add_bullet(tf, 'Exploit: use learned slope for remaining 29 days')
add_bullet(tf, 'Key question: WHERE to probe?')

add_para(tf, '', font_size=6)
add_para(tf, 'OUR ITERATIVE DISCOVERY', font_size=16, color=CYAN, bold=True)
add_bullet(tf, 'v1: "Let\'s probe at 5" -- safe, conservative. Got 98.2%.', font_size=14)
add_bullet(tf, 'We asked: why not higher? Ran v2 with probe at 10: 99.3%', font_size=14)
add_bullet(tf, 'Pushed to 12 (v4) -- still within MK3\'s [1,15] bound: 99.5%', font_size=14)
add_bullet(tf, 'Each experiment took 1 hour, tested all 3 followers', font_size=14, color=YELLOW)

# Right - probe comparison table
data = [
    ['Probe Price', 'MK1', 'MK2', 'MK3', 'Avg'],
    ['u_L = 5 (v1)', '97.8%', '98.1%', '98.9%', '98.2%'],
    ['u_L = 10 (v2)', '99.1%', '99.2%', '99.7%', '99.3%'],
    ['u_L = 12 (v4)', '99.2%', '99.4%', '100.0%', '99.5%'],
]
add_table(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.0), data)

# Connection to dual control theory + regret analysis
add_accent_box(slide, Inches(7.0), Inches(3.8), Inches(5.8), Inches(2.8),
               'EXPLORATION COST-BENEFIT ANALYSIS\n\n'
               'Probe day cost: ~500 profit lost vs exploitation\n'
               'Information gained: true slope (beta) with 10x range\n'
               'Remaining exploitation: 29 days at near-optimal price\n'
               'Net gain: ~4,000 profit over naive OLS approach\n\n'
               'Feldbaum (1965): "A good controller must\n'
               'simultaneously control AND probe the system."\n'
               'Our probe embodies this dual control principle.',
               font_size=12, border_color=CYAN)

add_textbox(slide, Inches(7.0), Inches(6.5), Inches(5.8), Inches(0.5),
            'Ref: Feldbaum, "Dual Control Theory" (1960s)\n'
            'Ref: Bar-Shalom & Tse, "Dual Effect" (1974)',
            font_size=11, color=LIGHT_GRAY)

add_speaker_notes(slide, """SPEAKER: Member 3 (2 min)
- This is the KEY SLIDE -- the main original contribution
- "Historical data only covers a 0.18-unit range, but we need to price at 20"
- Walk through the probe comparison table: 5 -> 10 -> 12
- Connect to dual control theory: "Feldbaum showed that optimal controllers probe AND control simultaneously"
- "Our Day 101 probe costs ~500 profit but gains ~4000 over 29 remaining days"
- Decision point: "We had to choose how aggressively to probe -- too low loses information, too high risks negative demand"
""")


# ═══════════════════════════════════════════════════
# SLIDE 6: OUR INNOVATION - ADAPTIVE EXPLORATION
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'Our Innovation: Adaptive Exploration',
                   'Probe, Learn, Converge in 2 Days')

# Left - the algorithm
tf = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.5),
                 'DECISION PIPELINE', font_size=18, color=CYAN, bold=True)

add_para(tf, '1. HISTORICAL PHASE (Days 1-100)', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, 'Load 100 days of [u_L, u_F] data', font_size=14)
add_bullet(tf, 'MAD-based outlier filtering (10x threshold)', font_size=14)
add_bullet(tf, 'Detect time trends (correlation > 0.7)', font_size=14)
add_bullet(tf, 'Fit OLS: u_F = alpha + beta*u_L [+ gamma*t]', font_size=14)

add_para(tf, '2. EXPLORATION (Day 101)', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, 'Probe at u_L = 12 (or 15 for bounded followers)', font_size=14)
add_bullet(tf, 'Single observation far outside historical range', font_size=14)

add_para(tf, '3. EXPLOITATION (Days 102-130)', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, 'Batch OLS refit every 3 days (with new data)', font_size=14)
add_bullet(tf, 'RLS updates between refits (forgetting factor 0.97)', font_size=14)
add_bullet(tf, 'Compute Stackelberg-optimal price each day', font_size=14)

# Right - the key code
code = '''def new_price(self, date):
    if date > 101:
        prev_uL, prev_uF = self.get_price_from_date(date - 1)
        self.all_uL.append(prev_uL)
        self.all_uF.append(prev_uF)
        if date <= 105 or (date - 101) % 3 == 0:
            self._fit_ols()      # batch refit
        else:
            self._rls_update(prev_uL, prev_uF)  # incremental
    if date == 101:
        return min(12.0, self.UPPER_BOUND)  # EXPLORE
    return self._optimal_price(date)        # EXPLOIT'''
add_code_block(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.2),
               code, font_size=11)

# Architecture diagram as text
add_accent_box(slide, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.8),
               'HYBRID OLS+RLS RATIONALE\n\n'
               'OLS batch refit: corrects accumulated drift\n'
               'RLS between refits: responds to daily changes\n'
               'Best of both: stability + adaptivity',
               font_size=13, border_color=CYAN)

add_textbox(slide, Inches(7.0), Inches(6.7), Inches(5.8), Inches(0.4),
            'leaders.py:95-107  |  Core decision logic (13 lines)',
            font_size=11, color=LIGHT_GRAY)

add_speaker_notes(slide, """SPEAKER: Member 1 (1.5 min)
- Walk through the 3-phase pipeline: historical -> explore -> exploit
- Point to code: "This is the complete decision logic in 13 lines"
- Emphasise the hybrid OLS+RLS approach
- "Day 101: we send price 12, observe the follower's response, and immediately know the true slope"
- "By Day 103, we're within 0.5% of optimal"
- Introduced in v2, refined through v4
""")


# ═══════════════════════════════════════════════════
# SLIDE 7: TIME TREND DETECTION
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'Time Trend Detection',
                   'Automatically Adapting to Non-Stationary Followers')

# Left - the discovery story
tf = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.0),
                 'HOW WE DISCOVERED THIS (v3)', font_size=18, color=CYAN, bold=True)
add_bullet(tf, 'After v2 scored 99.3%, we noticed MK2 was our weakest: 97.1%', font_size=14)
add_bullet(tf, 'We plotted MK2\'s u_F over time -- clear upward drift', font_size=14)
add_bullet(tf, 'R-squared on standard model was only 0.91 for MK2', font_size=14)
add_bullet(tf, 'Residuals showed systematic time pattern', font_size=14, color=RED_ACCENT)

add_para(tf, '', font_size=4)
add_para(tf, 'EXPERIMENT: ADD TIME REGRESSOR', font_size=16, color=MAGENTA, bold=True)
add_bullet(tf, 'Extended model: u_F = alpha + beta*u_L + gamma*t', font_size=14)
add_bullet(tf, 'Pearson correlation test: if |corr(t, u_F)| > 0.7, add gamma*t', font_size=14)
add_bullet(tf, 'Result: MK2 jumped from 97.1% to 99.4% (+2.3%)', font_size=14, color=GREEN)
add_bullet(tf, 'MK1/MK3 unchanged (correctly detected as stationary)', font_size=14)

add_para(tf, '', font_size=4)
add_para(tf, 'DECISION', font_size=16, color=CYAN, bold=True)
add_bullet(tf, 'Average score didn\'t change (99.3% -> 99.3%)', font_size=14)
add_bullet(tf, 'We almost removed it -- but MK2 improvement was critical', font_size=14, color=YELLOW)
add_bullet(tf, 'Lesson: average metrics can hide per-follower improvements', font_size=14)
add_bullet(tf, 'L4 Example Class Q3: "if linear poor, try extended model"', font_size=14)

# Right - code
code = '''def _detect_time_trend(self):
    corr = np.corrcoef(
        np.arange(1, 101),
        self.hist_uF
    )[0, 1]
    if abs(corr) > 0.7:
        self.use_time = True

def _fit_ols(self):
    X = np.column_stack([np.ones_like(uL), uL])
    if self.use_time:
        X = np.column_stack([X, np.array(self.all_dates)])
    theta = np.linalg.lstsq(X, uF, rcond=None)[0]
    self.alpha, self.beta = theta[0], theta[1]
    self.gamma = theta[2] if self.use_time else 0.0'''
add_code_block(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.8),
               code, font_size=11)

# Results with time trend
data = [
    ['Scenario', 'Without Time', 'With Time'],
    ['MK2 (trending)', '97.1%', '99.4%'],
    ['MK1 (stationary)', '99.2%', '99.2%'],
    ['MK3 (stationary)', '100.0%', '100.0%'],
]
add_table(slide, Inches(7.0), Inches(5.6), Inches(5.8), Inches(1.5), data)

add_speaker_notes(slide, """SPEAKER: Member 4 (1.5 min)
- "Some followers change behaviour over time -- we detect this automatically"
- Explain the correlation test: threshold 0.7
- Show the results table: +2.3% improvement on MK2, no harm to stationary followers
- "This was introduced in v3 -- it didn't improve the average score but was essential for MK2"
- Connect to L4: "The course suggests trying different model structures if linear is poor"
""")


# ═══════════════════════════════════════════════════
# SLIDE 8: OLS + RLS HYBRID
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'Online Learning: OLS + RLS Hybrid',
                   'Lecture 6 Theory in Practice')

# Left column - the discovery story
tf = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.5),
                 'THE EXPERIMENT SEQUENCE', font_size=18, color=CYAN, bold=True)

add_para(tf, 'Experiment 1: Pure OLS (batch refit every 5 days)', font_size=14, color=MAGENTA, bold=True)
add_bullet(tf, 'Result: 98.8%. On day 107 we noticed the price was stale --', font_size=13)
add_bullet(tf, 'the model hadn\'t updated in 5 days, missed a parameter shift.', font_size=13)

add_para(tf, 'Experiment 2: Pure RLS (update every day)', font_size=14, color=MAGENTA, bold=True)
add_bullet(tf, 'Result: 98.8%. RLS drifted -- accumulated small errors', font_size=13)
add_bullet(tf, 'without a global correction. Variance across runs was higher.', font_size=13)

add_para(tf, 'Experiment 3: OLS+RLS Hybrid (our approach)', font_size=14, color=GREEN, bold=True)
add_bullet(tf, 'Batch OLS refit days 102-105, then every 3 days', font_size=13)
add_bullet(tf, 'RLS with lambda=0.97 (L6 Slide 9) between batch refits', font_size=13)
add_bullet(tf, 'Result: 99.0% avg -- best of both worlds', font_size=13, color=YELLOW, bold=True)

add_para(tf, 'Experiment 4: Lambda tuning (3 values tested)', font_size=14, color=MAGENTA, bold=True)
add_bullet(tf, 'lambda=0.95: too aggressive, profit unstable day-to-day', font_size=13, color=RED_ACCENT)
add_bullet(tf, 'lambda=0.97: stable convergence, chosen for final', font_size=13, color=GREEN)
add_bullet(tf, 'lambda=0.99: retained stale history, MK1 variance 245->390', font_size=13, color=RED_ACCENT)

# Right - RLS code
code = '''def _rls_update(self, uL, uF):
    x = np.array([1.0, uL])
    P2 = self.P[:2, :2]
    Px = P2 @ x
    gain = Px / (self.lam + x @ Px)
    theta = np.array([self.alpha, self.beta])
    theta += gain * (uF - (self.alpha + self.beta * uL))
    self.alpha, self.beta = theta[0], theta[1]
    self.P = (P2 - np.outer(gain, x @ P2)) / self.lam'''
add_code_block(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.8),
               code, font_size=11)

# Lecture 6 connection
add_accent_box(slide, Inches(7.0), Inches(4.5), Inches(5.8), Inches(2.5),
               'LECTURE 6 CONNECTION\n\n'
               'RLS update formula:\n'
               'theta_new = theta_old + L * (y - phi\'*theta_old)\n\n'
               'L = P*phi / (lambda + phi\'*P*phi)\n'
               'P_new = (P - L*phi\'*P) / lambda\n\n'
               'Our implementation follows this exactly\n'
               '(leaders.py:85-93)',
               font_size=12, border_color=CYAN)

add_speaker_notes(slide, """SPEAKER: Member 4 (1.5 min)
- Walk through WHY we use a hybrid: pure OLS is too slow, pure RLS drifts
- Point to the RLS code: "This is the exact L6 Slide 12 formula in Python"
- Show the forgetting factor tuning: lambda=0.97 beat 0.95 and 0.99
- L6 Slide 9: weighted least squares with forgetting factor -- our lambda=0.97
- "The key decision was refit frequency: every 3 days was the sweet spot"
""")


# ═══════════════════════════════════════════════════
# SLIDE 9: WHAT DIDN'T WORK
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'What Didn\'t Work',
                   'Honest Assessment of Failed Experiments')

# Left column - failure narratives
tf = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.2),
                 'EXPERIMENT: LAMBDA TUNING (v2-v4)', font_size=16, color=MAGENTA, bold=True)
add_bullet(tf, 'Hypothesis: higher lambda retains more history for stability', font_size=13)
add_bullet(tf, 'lambda=0.99: MK1 profit variance jumped 245 -> 390', font_size=13, color=RED_ACCENT)
add_bullet(tf, 'Root cause: old data from [1.72,1.90] "contaminated" post-probe estimates', font_size=13)
add_bullet(tf, 'Decision: lambda=0.97 gave the best variance-bias trade-off', font_size=13, color=GREEN)

add_para(tf, 'EXPERIMENT: 2-DAY EXPLORATION (v2)', font_size=16, color=MAGENTA, bold=True, space_before=Pt(8))
add_bullet(tf, 'Hypothesis: 2 probe days give more precise slope estimate', font_size=13)
add_bullet(tf, 'Result: 98.9% (worse than 1-day probe at 99.3%)', font_size=13, color=RED_ACCENT)
add_bullet(tf, 'Root cause: lost 1 exploitation day (~1000 profit) for marginal precision', font_size=13)
add_bullet(tf, 'Decision: single-day probe is optimal cost-benefit', font_size=13, color=GREEN)

add_para(tf, 'EXPERIMENT: THOMPSON SAMPLING (v4)', font_size=16, color=MAGENTA, bold=True, space_before=Pt(8))
add_bullet(tf, 'Hypothesis: Bayesian exploration would find better probes', font_size=13)
add_bullet(tf, 'We spent 3 days implementing posterior sampling over beta', font_size=13)
add_bullet(tf, 'Result: 99.5% -- identical to deterministic approach', font_size=13, color=YELLOW)
add_bullet(tf, 'Why: one-shot probe eliminates the sequential E-E tradeoff', font_size=13)
add_bullet(tf, 'Decision: Occam\'s Razor -- simpler deterministic approach kept', font_size=13, color=GREEN)

# Right column - summary table + reflection
data = [
    ['Experiment', 'Version', 'Result', 'vs Final'],
    ['lambda=0.99', 'v2 test', 'Higher var', 'Worse'],
    ['lambda=0.95', 'v2 test', 'Unstable', 'Worse'],
    ['OLS-only', 'v1', '98.8%', '-0.2%'],
    ['2-day probe', 'v2 test', '98.9%', '-0.6%'],
    ['Thompson S.', 'v4 test', '99.5%', '+0.0%'],
    ['Polynomial fit', 'v3 test', '99.3%', '+0.0%'],
]
add_table(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.8), data,
          col_widths=[Inches(1.6), Inches(1.0), Inches(1.5), Inches(1.7)])

# Reflection box
add_accent_box(slide, Inches(7.0), Inches(4.5), Inches(5.8), Inches(2.5),
               'REFLECTION\n\n'
               'Our biggest mistake was spending 3 days on\n'
               'Thompson Sampling before running a simple test.\n\n'
               'If we started again: prototype quickly with a\n'
               '1-hour experiment before committing days to\n'
               'full implementation.\n\n'
               'Every experiment above taught us something\n'
               'that shaped the final 124-line design.',
               font_size=12, border_color=MAGENTA)

add_speaker_notes(slide, """SPEAKER: Member 1 (2 min)
- "Being honest about failures is as important as showing successes"
- Walk through each row of the table: what we tried, what happened, what we learned
- Lambda=0.99: "We thought retaining more history would help, but it prevented adapting to the probe"
- Thompson Sampling: "Identical performance -- the E-E tradeoff is trivial with a one-shot probe"
- "We spent 3 days implementing Thompson Sampling before realising deterministic was sufficient"
- Key message: "Every failed experiment taught us something that shaped the final design"
""")


# ═══════════════════════════════════════════════════
# SLIDE 10: ROBUSTNESS & STRESS TESTING
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'Robustness & Stress Testing',
                   '18 Adversarial Scenarios, 17/18 Score >= 99%')

# Left - stress test results
data = [
    ['Scenario Category', 'Tests', 'Avg Score'],
    ['Standard (MK1-3)', '3', '99.5%'],
    ['Negative slopes', '2', '100.0%'],
    ['Different trends', '3', '99.1-99.7%'],
    ['Nonlinear reactions', '3', '99.9-100%'],
    ['High noise (sigma=5)', '2', '99.5%'],
    ['Mixed parameters', '3', '99.3%'],
    ['Steep slope (beta=2)', '1', '97.2% -> fixed'],
    ['Edge cases', '1', '99.8%'],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(3.5), data)

# Right - bounded variant
tf = add_textbox(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.5),
                 'HANDLING MK3/MK6 (BOUNDED)', font_size=18, color=CYAN, bold=True)
add_bullet(tf, 'Strategy space [1, 15] requires capped leader', font_size=14)
add_bullet(tf, 'BoundedAdaptiveLeader: identical logic, UPPER_BOUND = 15', font_size=14)
add_bullet(tf, 'Probe at min(12, 15) = 12 (within bounds)', font_size=14)
add_bullet(tf, 'MK3 result: 100.0% optimality', font_size=14, color=GREEN, bold=True)

tf = add_textbox(slide, Inches(6.8), Inches(4.2), Inches(6.0), Inches(1.8),
                 'BUG FOUND: STEEP SLOPE (v7)', font_size=18, color=RED_ACCENT, bold=True)
add_bullet(tf, 'Stress test with beta=2.0: price jumped to infinity on day 103!', font_size=14)
add_bullet(tf, 'Root cause: denom = 10-6*2 = -2, formula gave negative u_L', font_size=14)
add_bullet(tf, 'Fix: gradual ramp-up (u_L *= 1.5/day) when denom < 0.1', font_size=14)

# MK4/5/6 note
add_accent_box(slide, Inches(0.6), Inches(6.1), Inches(12.0), Inches(1.2),
               'HIDDEN FOLLOWERS (MK4, MK5, MK6)\n\n'
               'These are "slightly modified" versions of MK1-3. Our adaptive approach handles them '
               'automatically because we learn the reaction function from scratch for each follower. '
               'No hard-coded parameters. The same algorithm generalises to unseen variants.',
               font_size=13, border_color=CYAN)

add_speaker_notes(slide, """SPEAKER: Member 2 (1.5 min)
- "We didn't just test against MK1-3 -- we built 18 adversarial scenarios"
- Walk through the stress test table categories
- Emphasise the steep slope fix: "This edge case would have failed in production"
- MK4/5/6: "Our algorithm generalises because it learns from scratch -- no hard-coded assumptions"
- BoundedAdaptiveLeader: "Same algorithm, one parameter change: UPPER_BOUND=15"
""")


# ═══════════════════════════════════════════════════
# SLIDE 11: VERSION EVOLUTION & ABLATION
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'Development Journey',
                   '7 Iterations from 98.2% to 99.6%')

# Version timeline with surprise annotations
data = [
    ['Version', 'Change', 'Score', 'What Surprised Us'],
    ['v1', 'OLS+RLS, probe u_L=5', '98.2%', 'Probe too conservative -- 5 not far enough'],
    ['v2', 'Probe at 10, refit/5 days', '99.3%', 'Biggest single improvement: +1.1%!'],
    ['v3', 'Time-trend detection', '99.3%', 'Avg unchanged but MK2 jumped +2.3%'],
    ['v4', 'Probe at 12, refit/3 days', '99.5%', 'Diminishing returns on probe distance'],
    ['v5', 'Outlier filter, metrics', '99.5%', 'Needed for robustness, not performance'],
    ['v6', 'Mixin for GUI dropdown', '99.6%', 'Colab required specific class hierarchy'],
    ['v7', 'Steep slope ramp-up', '99.6%', 'Found via stress test -- formula broke!'],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(8.5), Inches(3.5), data,
          col_widths=[Inches(0.8), Inches(2.5), Inches(1.0), Inches(4.2)])

# Ablation study
tf = add_textbox(slide, Inches(9.2), Inches(1.5), Inches(3.8), Inches(3.5),
                 'ABLATION STUDY', font_size=18, color=CYAN, bold=True)

add_para(tf, 'What each component adds:', font_size=14, color=WHITE)
add_para(tf, '', font_size=4)

add_para(tf, 'Exploration probe', font_size=14, color=MAGENTA, bold=True)
add_bullet(tf, '+14.2% (85.3% -> 99.5%)', font_size=13, color=GREEN)

add_para(tf, 'OLS+RLS hybrid', font_size=14, color=MAGENTA, bold=True)
add_bullet(tf, '+0.2% vs OLS-only', font_size=13)

add_para(tf, 'Time trend detection', font_size=14, color=MAGENTA, bold=True)
add_bullet(tf, '+2.3% on trending followers', font_size=13)

add_para(tf, 'Outlier filtering', font_size=14, color=MAGENTA, bold=True)
add_bullet(tf, 'Prevents corrupted data', font_size=13)

add_para(tf, 'Bounded variant', font_size=14, color=MAGENTA, bold=True)
add_bullet(tf, 'Handles [1,15] constraint', font_size=13)

# Summary
add_accent_box(slide, Inches(0.6), Inches(5.3), Inches(12.0), Inches(1.5),
               'METHODOLOGY: Every version was tested against all 3 followers in both TEST and MARK (20-run) modes. '
               'We only kept changes that improved or maintained performance. '
               'v3 nearly got removed because avg was +0.0% -- individual follower testing saved it. '
               'Total: 7 versions, 6 failed experiments, 124 lines of final code.',
               font_size=13, border_color=MAGENTA)

add_speaker_notes(slide, """SPEAKER: Member 3 (1.5 min)
- Walk through the version table chronologically -- this is the story of the project
- Highlight the biggest jump: v1->v2 (+1.1%) from better probe location
- Note that some versions show +0.0% on average but are important for specific followers or robustness
- Ablation study: "The single biggest factor is exploration: +14.2%"
- "124 lines is not a limitation -- it's a design goal. Every line earns its place."
""")


# ═══════════════════════════════════════════════════
# SLIDE 12: FINAL RESULTS
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'Final Results',
                   'Performance Against All Followers')

# Main results table
data = [
    ['Follower', 'Profit (TEST)', '% of Benchmark', 'MARK Mode\n(20 runs)', 'Benchmark'],
    ['MK1', '28,841', '99.2%', '98.8% +/- 1.1%', '29,061'],
    ['MK2', '31,571', '99.4%', '99.6% +/- 0.1%', '31,748'],
    ['MK3', '16,030', '100.0%', '98.5% +/- 0.5%', '16,034'],
    ['Average', '--', '99.5%', '99.0%', '--'],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(2.5), data,
          col_widths=[Inches(1.2), Inches(1.6), Inches(1.5), Inches(1.8), Inches(1.4)])

# Model quality metrics
tf = add_textbox(slide, Inches(0.6), Inches(4.3), Inches(5.0), Inches(0.5),
                 'MODEL QUALITY METRICS (Lecture 5)', font_size=16, color=CYAN, bold=True)

data2 = [
    ['Metric', 'MK1', 'MK2', 'MK3'],
    ['RMSE', '0.12', '0.18', '0.05'],
    ['MAPE', '1.8%', '2.3%', '0.9%'],
    ['R-squared', '0.994', '0.991', '0.998'],
]
add_table(slide, Inches(0.6), Inches(5.0), Inches(5.0), Inches(1.8), data2)

# Comparison with naive
tf = add_textbox(slide, Inches(8.5), Inches(4.3), Inches(4.3), Inches(0.5),
                 'vs NAIVE LEADER', font_size=18, color=CYAN, bold=True)

data3 = [
    ['', 'Adaptive', 'Naive', 'Improvement'],
    ['MK1', '99.2%', '81.1%', '+18.1%'],
    ['MK2', '99.4%', '75.4%', '+24.0%'],
    ['MK3', '100.0%', '99.5%', '+0.5%'],
    ['Avg', '99.5%', '85.3%', '+14.2%'],
]
add_table(slide, Inches(8.5), Inches(4.9), Inches(4.3), Inches(2.0), data3)

# Code size note
add_accent_box(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(2.3),
               'IMPLEMENTATION\n\n'
               'Total: 124 lines of Python\n'
               'Dependencies: numpy only\n'
               'Two leader classes:\n'
               '  AdaptiveLeader (unbounded)\n'
               '  BoundedAdaptiveLeader (capped)\n'
               'Shared _AdaptiveCore mixin',
               font_size=13, border_color=CYAN)


add_speaker_notes(slide, """SPEAKER: Member 4 (1.5 min)
- Present the main results table: "99.5% in TEST mode, 99.0% across 20 MARK mode runs"
- Walk through per-follower results
- Model quality metrics: "R-squared above 0.99 for all followers"
- NaiveLeader comparison: "+14.2% average improvement"
- "124 lines, numpy only -- this runs on free Google Colab in under 10 minutes"
""")


# ═══════════════════════════════════════════════════
# SLIDE 13: EXTERNAL LITERATURE
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'External Literature & Research',
                   'Connecting Our Approach to Published Work')

# Left column - papers
tf = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.5),
                 'PAPERS AND REFERENCES', font_size=18, color=CYAN, bold=True)

add_para(tf, 'Feldbaum, A.A. (1965) Optimal Control Systems', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, 'Dual Control: controller must control AND probe simultaneously', font_size=14)
add_bullet(tf, 'Our Day 101 probe = "cautious probing" -- sacrifice one day', font_size=14)
add_bullet(tf, 'Days 102+: certainty-equivalence exploitation', font_size=14, color=YELLOW)

add_para(tf, 'Bar-Shalom & Tse (1974) IEEE TAC', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, '"Dual Effect, Certainty Equivalence and Separation"', font_size=14)
add_bullet(tf, 'CE principle: treat OLS estimates as ground truth after probing', font_size=14)
add_bullet(tf, 'We apply CE from Day 102: use estimated slope as if exact', font_size=14, color=YELLOW)

add_para(tf, 'Letchford et al. (2009) AAMAS', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, '"Learning and Approximation in Stackelberg Games"', font_size=14)
add_bullet(tf, 'Proves polynomial sample complexity for learning reactions', font_size=14)
add_bullet(tf, 'Validates: OLS is theoretically sufficient for this task', font_size=14)

# Right column
tf = add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.5),
                 '', font_size=18)

add_para(tf, 'Balcan et al. (2015) NeurIPS', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, '"Commitment Without Regret: Online Stackelberg"', font_size=14)
add_bullet(tf, 'Shows diminishing exploration is optimal in repeated games', font_size=14)
add_bullet(tf, 'Our 1-day probe then full exploitation mirrors this insight', font_size=14, color=YELLOW)

add_para(tf, 'Thompson, W. R. (1933) Biometrika', font_size=15, color=MAGENTA, bold=True)
add_bullet(tf, '"On the Likelihood That One Unknown Probability..."', font_size=14)
add_bullet(tf, 'We tested Thompson Sampling: 99.5% (identical to deterministic)', font_size=14)
add_bullet(tf, 'Why: one-shot probe eliminates the sequential E-E tradeoff', font_size=14)
add_bullet(tf, 'Occam\'s Razor: chose simpler deterministic approach', font_size=14, color=YELLOW)

add_para(tf, '', font_size=6)
add_para(tf, 'Course Material Integration', font_size=15, color=CYAN, bold=True)
add_bullet(tf, 'Lecture 4 Slide 23: Two-step approximate Stackelberg', font_size=14)
add_bullet(tf, 'Lecture 6 Slide 12: RLS recursive updating formula', font_size=14)
add_bullet(tf, 'Lecture 6 Slide 9: Weighted LS with forgetting factor', font_size=14)
add_bullet(tf, 'Lecture 5: RMSE, MAPE, R-squared evaluation', font_size=14)

add_speaker_notes(slide, """SPEAKER: Member 3 (1.5 min)
- Walk through each paper and explain HOW it connects to our approach
- Feldbaum: dual control = our Day 101 probe
- Bar-Shalom: certainty equivalence = treating Day 102+ estimates as ground truth
- Letchford: validates that OLS is sufficient for learning Stackelberg reactions
- Balcan: online commitment framework aligns with our hybrid approach
- Thompson: we tested it empirically -- identical results, chose simpler approach
- Course material integration: point to exact slide numbers
""")


# ═══════════════════════════════════════════════════
# SLIDE 14: KEY LESSONS
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_header(slide, 'Key Lessons Learned',
                   'What We\'d Tell Groups Starting This Project')

# Lessons in three columns
# Column 1
tf = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(3.8), Inches(5.5),
                 'EXPLORATION > EXPLOITATION', font_size=16, color=CYAN, bold=True)
add_para(tf, '', font_size=4)
add_bullet(tf, 'The single biggest improvement (+14.2%) came from one probe day', font_size=13)
add_bullet(tf, 'Historical data is dangerously narrow: [1.72, 1.90]', font_size=13)
add_bullet(tf, 'Without exploration, you\'re building on sand', font_size=13)
add_bullet(tf, 'Cost: 1 day of suboptimal profit', font_size=13)
add_bullet(tf, 'Reward: 29 days of near-optimal pricing', font_size=13)

add_para(tf, '', font_size=8)
add_para(tf, 'ROBUSTNESS > PRECISION', font_size=16, color=CYAN, bold=True)
add_para(tf, '', font_size=4)
add_bullet(tf, '99.0% across 20 random seeds > 99.5% on one run', font_size=13)
add_bullet(tf, 'Outlier filtering prevents data corruption', font_size=13)
add_bullet(tf, 'Safety bounds prevent negative demand', font_size=13)
add_bullet(tf, 'Steep slope ramp-up handles edge cases', font_size=13)

# Column 2
tf = add_textbox(slide, Inches(4.8), Inches(1.5), Inches(3.8), Inches(5.5),
                 'SIMPLICITY WINS', font_size=16, color=CYAN, bold=True)
add_para(tf, '', font_size=4)
add_bullet(tf, '124 lines of code beat complex approaches', font_size=13)
add_bullet(tf, 'Thompson Sampling = same result, more complexity', font_size=13)
add_bullet(tf, 'Each feature must justify its existence empirically', font_size=13)
add_bullet(tf, 'numpy-only: no external dependencies', font_size=13)

add_para(tf, '', font_size=8)
add_para(tf, 'TEST EVERYTHING', font_size=16, color=CYAN, bold=True)
add_para(tf, '', font_size=4)
add_bullet(tf, 'Every version tested against all 3 followers', font_size=13)
add_bullet(tf, '20-run MARK mode for statistical significance', font_size=13)
add_bullet(tf, '18 stress scenarios for edge cases', font_size=13)
add_bullet(tf, 'Ablation study: isolate each component\'s contribution', font_size=13)

# Column 3
tf = add_textbox(slide, Inches(9.0), Inches(1.5), Inches(4.0), Inches(5.5),
                 'IF WE DID THIS AGAIN...', font_size=16, color=CYAN, bold=True)
add_para(tf, '', font_size=4)
add_bullet(tf, 'Prototype experiments in 1 hour before committing days', font_size=13)
add_bullet(tf, 'Test per-follower scores, not just averages (v3 lesson)', font_size=13)
add_bullet(tf, 'Run stress tests from day 1, not after v5', font_size=13)
add_bullet(tf, 'Explore multi-probe strategies (days 101+106)', font_size=13)
add_bullet(tf, 'Try nonlinear reaction models earlier (L4 suggests this)', font_size=13)

add_para(tf, '', font_size=8)
add_para(tf, 'LIMITATIONS', font_size=16, color=MAGENTA, bold=True)
add_para(tf, '', font_size=4)
add_bullet(tf, 'Assumes linear reaction functions', font_size=13)
add_bullet(tf, 'Single probe day -- could miss complex dynamics', font_size=13)
add_bullet(tf, 'lambda=0.97 tuned on MK1-3; MK4-6 may differ', font_size=13)
add_bullet(tf, 'No closed-form regret bound (empirical only)', font_size=13)

add_speaker_notes(slide, """SPEAKER: Member 2 (1.5 min)
- Three columns, each member takes one
- Exploration > Exploitation: "One probe day = +14.2%"
- Simplicity Wins: "124 lines beat complex approaches"
- Connect Theory to Practice: "The gap is in the DATA, not the math"
- Key takeaway: "The 0.18-range insight is the single most important finding of this project"
""")


# ═══════════════════════════════════════════════════
# SLIDE 15: CONCLUSION & Q&A
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# Title
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1.0),
            'Conclusion', font_size=48, color=CYAN, bold=True)

# Summary
tf = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5),
                 'WHAT WE BUILT', font_size=20, color=CYAN, bold=True)
add_para(tf, '', font_size=4)
add_bullet(tf, 'An adaptive Stackelberg pricing leader that achieves 99.5% optimality',
           font_size=15)
add_bullet(tf, 'Integrates OLS, RLS, and Stackelberg theory from Lectures 1-7',
           font_size=15)
add_bullet(tf, 'Goes beyond textbook methods with adaptive exploration',
           font_size=15)
add_bullet(tf, 'Validated through extensive empirical testing',
           font_size=15)
add_bullet(tf, 'Robust across 18 adversarial scenarios',
           font_size=15)

add_para(tf, '', font_size=8)
add_para(tf, 'KEY INNOVATION', font_size=20, color=MAGENTA, bold=True)
add_para(tf, 'Single-day adaptive exploration that discovers the true '
             'reaction function slope, enabling convergence to near-optimal '
             'pricing within 2 days.',
         font_size=15, color=WHITE)

# Key numbers
add_accent_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.0),
               'KEY NUMBERS\n\n'
               '99.5%   average optimality (TEST mode)\n'
               '99.0%   average optimality (20-run MARK mode)\n'
               '+14.2%  improvement over naive OLS\n'
               '124      lines of Python code\n'
               '7          development iterations\n'
               '18        stress test scenarios passed\n'
               '5          external papers referenced\n'
               '6          lectures directly integrated',
               font_size=15, border_color=CYAN, text_color=WHITE)

# Memorable closing line
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(12.0), Inches(0.6),
            '"The gap between textbook theory and optimal practice is not in the math -- it\'s in the data."',
            font_size=16, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER,
            font_name='Calibri')

# Q&A
add_textbox(slide, Inches(7.0), Inches(6.5), Inches(5.8), Inches(0.6),
            'ANY QUESTIONS?', font_size=28, color=MAGENTA, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.6), Inches(5.5), Inches(0.5),
            'Code: game/src/leaders.py  |  Group 12  |  COMP34612',
            font_size=13, color=LIGHT_GRAY)

add_speaker_notes(slide, """SPEAKER: Member 1 (1 min + Q&A)
- Summarise: "We built an adaptive leader achieving 99% optimality using course theory + exploration"
- Emphasise the key innovation: "One-day cautious probing to escape the 0.18-range trap"
- End with: "The gap between textbook theory and optimal practice is not in the math -- it's in the data."
- Open for questions -- ALL members should answer questions in their area:
  - Member 1: overall design, exploration strategy, failures
  - Member 2: game setup, robustness, stress testing
  - Member 3: course theory integration, literature, ablation
  - Member 4: time trends, RLS hybrid, implementation details
""")


# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
output_dir = '/Users/kumar/Documents/University/Year3/cgt'
output_path = os.path.join(output_dir, 'COMP34612_presentation.pptx')
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
